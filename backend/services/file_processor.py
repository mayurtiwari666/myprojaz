import boto3
import pypdf
import io
import os
from docx import Document
from pptx import Presentation
from backend.config import settings

s3_client = boto3.client('s3', region_name=settings.AWS_REGION)

# Initialize Textract with Specific Credentials (Secondary Account)
textract_client = boto3.client(
    'textract',
    region_name=settings.AWS_TEXTRACT_REGION,
    aws_access_key_id=settings.AWS_TEXTRACT_ACCESS_KEY_ID,
    aws_secret_access_key=settings.AWS_TEXTRACT_SECRET_ACCESS_KEY
)

def extract_text_from_s3(file_key: str) -> str:
    """
    Downloads a file from S3 and extracts its text based on extension.
    Supports: .pdf (with Textract fallback), .docx, .pptx, .txt
    """
    try:
        ext = os.path.splitext(file_key)[1].lower()
        print(f"Extraction started for {file_key} ({ext})")
        
        response = s3_client.get_object(Bucket=settings.S3_BUCKET_NAME, Key=file_key)
        file_content = response['Body'].read()
        file_stream = io.BytesIO(file_content)
        
        if ext == '.pdf':
            print(f"PDF detected: {file_key}. Using AWS Textract (Policy: Always Use Textract).")
            # Force Textract call directly using BYTES to avoid cross-account S3 issues
            return _extract_with_textract(file_content)
        elif ext in ['.docx', '.doc']:
            return _extract_from_docx(file_stream)
        elif ext in ['.pptx', '.ppt']:
            return _extract_from_pptx(file_stream)
        elif ext in ['.txt', '.md']:
            return file_content.decode('utf-8', errors='ignore')
        elif ext in ['.jpg', '.jpeg', '.png']:
            print("Image file detected. Skipping text extraction as per budget policy.")
            return "" # Return empty string for images (Metadata only)
        else:
            print(f"Unsupported file type for extraction: {ext}")
            return ""
            
    except Exception as e:
        print(f"Error extracting text from {file_key}: {e}")
        raise e

def _extract_from_pdf(file_stream, file_key) -> str:
    """
    Hybrid extraction: pypdf first, Textract fallback if empty.
    """
    text = ""
    try:
        reader = pypdf.PdfReader(file_stream)
        for page in reader.pages:
            extract = page.extract_text()
            if extract:
                text += extract + "\n"
    except Exception as e:
        print(f"pypdf failed: {e}")
    
    # Check if scanned (empty or very short text)
    # Debug: Check what pypdf actually found
    print(f"pypdf extracted {len(text)} chars. Preview: {text[:200]!r}")

    # Check 1: Length (Too short?)
    if not text or len(text.strip()) < 200:
        print(f"PDF text length ({len(text.strip())}) < 200. Marking as scanned.")
        return _extract_with_textract(settings.S3_BUCKET_NAME, file_key)

    # Check 2: Density (Garbage/Symbols?)
    # Count alphanumeric chars (a-z, 0-9)
    alphanumeric_count = sum(c.isalnum() for c in text)
    density = alphanumeric_count / len(text) if len(text) > 0 else 0
    
    if density < 0.5:
        print(f"PDF text density ({density:.2f}) < 0.5. likely garbage OCR or symbols. Falling back to AWS Textract...")
        return _extract_with_textract(settings.S3_BUCKET_NAME, file_key)
        
    return text

def _extract_with_textract(file_bytes) -> str:
    """
    Uses AWS Textract to detect text in a document using raw bytes.
    This avoids cross-account S3 permission issues.
    """
    try:
        response = textract_client.detect_document_text(
            Document={'Bytes': file_bytes}
        )
        blocks = response.get('Blocks', [])
        text = ""
        for block in blocks:
            if block['BlockType'] == 'LINE':
                text += block['Text'] + "\n"
        print(f"Textract successfully extracted {len(text)} chars. Preview: {text[:200]!r}")
        return text
    except Exception as e:
        print(f"Textract failed: {e}")
        return ""

def _extract_from_docx(file_stream) -> str:
    try:
        doc = Document(file_stream)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Docx extraction failed: {e}")
        return ""

def _extract_from_pptx(file_stream) -> str:
    try:
        prs = Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"PPTX extraction failed: {e}")
        return ""
